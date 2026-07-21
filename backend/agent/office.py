"""Dedicated local office-file builders and structural inspectors (WB-243).

All writers use a sibling temporary file and os.replace, so a failed build or
validation never leaves a partial deliverable at the requested sandbox path.
"""
from __future__ import annotations

from html import escape
import os
from pathlib import Path
import uuid
from typing import Any, Callable

from agent.sandbox import relpath, resolve_in_sandbox


class OfficeValidationError(ValueError):
    pass


def _path(path: str, suffix: str) -> Path:
    target = resolve_in_sandbox(path)
    if target.suffix.lower() != suffix:
        raise OfficeValidationError(f"文件扩展名必须是 {suffix}")
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def _atomic_build(
    path: str, suffix: str, builder: Callable[[Path], None],
    inspector: Callable[[Path], dict[str, Any]],
) -> tuple[str, dict[str, Any]]:
    target = _path(path, suffix)
    temp = target.with_name(f".{target.stem}.{uuid.uuid4().hex}.tmp{suffix}")
    try:
        builder(temp)
        validation = inspector(temp)
        if not validation.get("valid"):
            raise OfficeValidationError(str(validation.get("errors") or "格式校验失败"))
        os.replace(temp, target)
        return relpath(target), validation
    finally:
        if temp.exists():
            temp.unlink()


def inspect_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    headings = [p.text for p in doc.paragraphs if p.style and p.style.name.startswith("Heading")]
    table_rows = sum(len(table.rows) for table in doc.tables)
    return {
        "valid": True, "format": "docx", "paragraphs": len(paragraphs),
        "headings": len(headings), "tables": len(doc.tables), "table_rows": table_rows,
        "visual_rendered": False,
    }


def create_docx(args: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from docx import Document
    from docx.shared import Pt

    def build(temp: Path) -> None:
        doc = Document()
        normal = doc.styles["Normal"]
        normal.font.name = str(args.get("font") or "Arial")[:80]
        normal.font.size = Pt(float(args.get("font_size") or 11))
        title = str(args.get("title") or "").strip()
        if title:
            doc.add_heading(title, level=0)
        for section in (args.get("sections") or [])[:100]:
            heading = str(section.get("heading") or "").strip()
            if heading:
                doc.add_heading(heading, level=max(1, min(int(section.get("level") or 1), 3)))
            for paragraph in (section.get("paragraphs") or [])[:500]:
                doc.add_paragraph(str(paragraph))
        for table_data in (args.get("tables") or [])[:30]:
            rows = (table_data.get("rows") or [])[:1000]
            if not rows:
                continue
            width = max(len(row) for row in rows)
            table = doc.add_table(rows=len(rows), cols=width)
            table.style = "Table Grid"
            for r_idx, row in enumerate(rows):
                for c_idx, value in enumerate(row[:width]):
                    table.cell(r_idx, c_idx).text = str(value)
        doc.save(temp)

    return _atomic_build(str(args["path"]), ".docx", build, inspect_docx)


def inspect_xlsx(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    book = load_workbook(path, data_only=False)
    formulas = 0
    charts = 0
    non_empty = 0
    for sheet in book.worksheets:
        charts += len(sheet._charts)
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    non_empty += 1
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    formulas += 1
    return {
        "valid": True, "format": "xlsx", "sheets": len(book.worksheets),
        "sheet_names": book.sheetnames, "non_empty_cells": non_empty,
        "formulas": formulas, "charts": charts, "recalculation_on_open": True,
        "visual_rendered": False,
    }


def create_xlsx(args: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.styles import Font, PatternFill

    def build(temp: Path) -> None:
        book = Workbook()
        specs = (args.get("sheets") or [])[:20]
        if not specs:
            specs = [{"name": "Sheet1", "rows": []}]
        for idx, spec in enumerate(specs):
            sheet = book.active if idx == 0 else book.create_sheet()
            raw_name = str(spec.get("name") or f"Sheet{idx + 1}")
            sheet.title = raw_name[:31].replace("/", "-").replace("\\", "-")
            rows = (spec.get("rows") or [])[:5000]
            for row in rows:
                sheet.append(list(row)[:100])
            if rows:
                for cell in sheet[1]:
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill("solid", fgColor="D9EAF7")
            for formula in (spec.get("formulas") or [])[:1000]:
                cell = str(formula.get("cell") or "").upper()
                value = str(formula.get("formula") or "")
                if cell and value.startswith("="):
                    sheet[cell] = value
            for col in sheet.columns:
                letter = col[0].column_letter
                width = min(60, max(10, max((len(str(cell.value or "")) for cell in col), default=0) + 2))
                sheet.column_dimensions[letter].width = width
            chart_spec = spec.get("chart")
            if chart_spec:
                chart_type = str(chart_spec.get("type") or "bar")
                chart = {"line": LineChart, "pie": PieChart}.get(chart_type, BarChart)()
                min_row = max(1, int(chart_spec.get("min_row") or 1))
                max_row = max(min_row, int(chart_spec.get("max_row") or sheet.max_row))
                min_col = max(1, int(chart_spec.get("data_min_col") or 2))
                max_col = max(min_col, int(chart_spec.get("data_max_col") or min_col))
                data = Reference(sheet, min_col=min_col, max_col=max_col, min_row=min_row, max_row=max_row)
                chart.add_data(data, titles_from_data=min_row == 1)
                category_col = max(1, int(chart_spec.get("categories_col") or 1))
                categories = Reference(sheet, min_col=category_col, min_row=min_row + (1 if min_row == 1 else 0), max_row=max_row)
                chart.set_categories(categories)
                chart.title = str(chart_spec.get("title") or "")[:120]
                sheet.add_chart(chart, str(chart_spec.get("anchor") or "G2"))
        book.calculation.fullCalcOnLoad = True
        book.calculation.forceFullCalc = True
        book.save(temp)

    return _atomic_build(str(args["path"]), ".xlsx", build, inspect_xlsx)


def inspect_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    deck = Presentation(path)
    violations: list[dict[str, Any]] = []
    text_shapes = 0
    for slide_no, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > deck.slide_width or shape.top + shape.height > deck.slide_height:
                violations.append({"slide": slide_no, "shape": getattr(shape, "name", "shape")})
    return {
        "valid": not violations, "format": "pptx", "slides": len(deck.slides),
        "text_shapes": text_shapes, "bounds_violations": violations,
        "visual_rendered": False,
    }


def create_pptx(args: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation
    from pptx.util import Pt

    def build(temp: Path) -> None:
        deck = Presentation()
        slides = (args.get("slides") or [])[:100]
        if not slides:
            slides = [{"title": str(args.get("title") or "Presentation"), "bullets": []}]
        for idx, spec in enumerate(slides):
            layout = deck.slide_layouts[0] if idx == 0 else deck.slide_layouts[1]
            slide = deck.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = str(spec.get("title") or "")[:300]
            body = next((shape for shape in slide.placeholders if shape != slide.shapes.title and getattr(shape, "has_text_frame", False)), None)
            bullets = (spec.get("bullets") or [])[:30]
            if body is not None and bullets:
                frame = body.text_frame
                frame.clear()
                for b_idx, bullet in enumerate(bullets):
                    paragraph = frame.paragraphs[0] if b_idx == 0 else frame.add_paragraph()
                    paragraph.text = str(bullet)[:1000]
                    paragraph.font.size = Pt(22)
                    paragraph.level = 0
        deck.save(temp)

    return _atomic_build(str(args["path"]), ".pptx", build, inspect_pptx)


def inspect_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    extracted = "".join((page.extract_text() or "") for page in reader.pages)
    return {
        "valid": len(reader.pages) > 0, "format": "pdf", "pages": len(reader.pages),
        "extractable_characters": len(extracted.strip()), "encrypted": bool(reader.is_encrypted),
        "visual_rendered": False,
    }


def create_pdf(args: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    def build(temp: Path) -> None:
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            font = "STSong-Light"
        except Exception:  # pragma: no cover - ReportLab builds may omit CID maps
            font = "Helvetica"
        styles = getSampleStyleSheet()
        body = ParagraphStyle("AgentMateBody", parent=styles["BodyText"], fontName=font, leading=17)
        heading = ParagraphStyle("AgentMateHeading", parent=styles["Heading1"], fontName=font)
        story = []
        title = str(args.get("title") or "").strip()
        if title:
            story.extend([Paragraph(escape(title), heading), Spacer(1, 12)])
        for paragraph in (args.get("paragraphs") or [])[:1000]:
            story.extend([Paragraph(escape(str(paragraph)).replace("\n", "<br/>"), body), Spacer(1, 8)])
        for table_data in (args.get("tables") or [])[:30]:
            rows = [[str(value) for value in row[:30]] for row in (table_data.get("rows") or [])[:1000]]
            if not rows:
                continue
            table = Table(rows, repeatRows=1)
            table.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), font), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.extend([table, Spacer(1, 12)])
        if not story:
            story.append(Paragraph(" ", body))
        SimpleDocTemplate(str(temp), pagesize=A4, title=title).build(story)

    return _atomic_build(str(args["path"]), ".pdf", build, inspect_pdf)


def inspect_office_file(path: str) -> dict[str, Any]:
    target = resolve_in_sandbox(path)
    if not target.is_file():
        raise FileNotFoundError(path)
    inspector = {
        ".docx": inspect_docx, ".xlsx": inspect_xlsx,
        ".pptx": inspect_pptx, ".pdf": inspect_pdf,
    }.get(target.suffix.lower())
    if not inspector:
        raise OfficeValidationError("只支持 .docx/.xlsx/.pptx/.pdf")
    return inspector(target)
